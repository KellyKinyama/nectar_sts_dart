"""Build the STS in-house reference proposal deck.

Run: python docs/build_pptx.py
Output: docs/sts-proposal.pptx  (and director-briefing.pptx)

Updated 2026-06-11 to reflect the broader work shipped:
  * nectar_sts_dart — pure-Dart implementation of the licensed STS
    primitives (DKGA-02 / DKGA-04 / EA07 STA / EA09 DEA / EA11 MISTY1)
    plus the full token-class hierarchy (Class 0/1/2 incl. paired and
    4-section MISTY1 KCT).
  * NectarAPI-compatible HTTP server (shelf) wrapping VirtualHsm.
  * nectar_virtual_meter — Flutter customer-side meter simulator
    (Windows / Android / Web / Linux) on top of the same library.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "sts-proposal.pptx"
OUT_ALT = Path(__file__).resolve().parent / "director-briefing.pptx"

# Brand-ish palette (neutral, projector-safe)
NAVY = RGBColor(0x0B, 0x2C, 0x4A)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
ACCENT = RGBColor(0xE0, 0x6C, 0x00)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
TEXT = RGBColor(0x1B, 0x1B, 0x1B)
MUTED = RGBColor(0x55, 0x5B, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align="left", font="Calibri"):
    from pptx.enum.text import PP_ALIGN
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {
            "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[align]
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.06), TEAL)
    add_text(slide, Inches(0.5), Inches(0.18), SW - Inches(1.0),
             Inches(0.6), title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), SW - Inches(1.0),
                 Inches(0.35), subtitle, size=14, color=LIGHT)


def footer(slide, page, total):
    add_text(slide, Inches(0.5), SH - Inches(0.4),
             Inches(8), Inches(0.3),
             "In-House STS Reference  ·  Proposal  (rev 11-Jun-2026)",
             size=10, color=MUTED)
    add_text(slide, SW - Inches(2.0), SH - Inches(0.4),
             Inches(1.5), Inches(0.3),
             f"{page} / {total}", size=10, color=MUTED, align="right")


def bullets(slide, x, y, w, h, items, *, size=18, color=TEXT, bold_first=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"\u2022  {item}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold_first and i == 0
    return box


def table(slide, x, y, w, h, rows, *, header_fill=NAVY,
          header_color=WHITE, body_size=14, header_size=14,
          col_widths=None):
    n_rows = len(rows)
    n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            if r_idx == 0:
                run.font.bold = True
                run.font.size = Pt(header_size)
                run.font.color.rgb = header_color
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                run.font.size = Pt(body_size)
                run.font.color.rgb = TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 else LIGHT
    return tbl


# --- Slide 1: Title ----------------------------------------------------
def s_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(3.1), SW, Inches(0.08), ACCENT)
    add_text(s, Inches(0.8), Inches(1.2), SW - Inches(1.6), Inches(1.4),
             "In-House STS Token Engine",
             size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.2), SW - Inches(1.6), Inches(1.0),
             "Vendor-side algorithm port + customer-side meter app \u2014 "
             "working end-to-end. Pending only STS certification.",
             size=22, color=LIGHT)
    add_text(s, Inches(0.8), Inches(3.6), SW - Inches(1.6), Inches(0.4),
             "Internal R&D track  \u00b7  Update + Proposal",
             size=18, color=ACCENT, bold=True)
    add_text(s, Inches(0.8), Inches(5.6), SW - Inches(1.6), Inches(0.4),
             "Kelly Kinyama  \u00b7  Thursday, 11 June 2026",
             size=16, color=WHITE)
    return s


# --- Slide 2: The ask --------------------------------------------------
def s_ask():
    s = prs.slides.add_slide(BLANK)
    header(s, "The ask, in one sentence")
    add_rect(s, Inches(1.0), Inches(1.8), SW - Inches(2.0), Inches(3.8),
             LIGHT)
    add_text(s, Inches(1.4), Inches(2.0), SW - Inches(2.8), Inches(0.6),
             "Sponsor the certification track for the in-house STS engine.",
             size=24, bold=True, color=NAVY)
    add_text(s, Inches(1.4), Inches(2.9), SW - Inches(2.8), Inches(0.6),
             "~$2,500 to buy the official IEC 62055-41 ed3.0 standard "
             "+ STS Association test vectors.",
             size=18, color=NAVY)
    add_text(s, Inches(1.4), Inches(3.7), SW - Inches(2.8), Inches(0.4),
             "Engineering is already done. HSM, certification fees and "
             "shadow-run effort are tracked separately.",
             size=14, color=MUTED)
    add_text(s, Inches(1.4), Inches(4.3), SW - Inches(2.8), Inches(0.6),
             "Outcome: a source-code STS stack \u2014 vendor side + "
             "customer side \u2014 that we own and control end to end.",
             size=20, color=NAVY)
    add_text(s, Inches(1.0), Inches(6.0), SW - Inches(2.0), Inches(0.6),
             "Lower the cost of every alternative \u2014 including this vendor.",
             size=22, bold=True, color=ACCENT, align="center")
    return s


# --- Slide 3: The numbers ----------------------------------------------
def s_numbers():
    s = prs.slides.add_slide(BLANK)
    header(s, "The numbers",
           "Vendor stays. We're buying certification, not a rewrite.")
    rows = [
        ["Item", "Today", "With alternative"],
        ["Vendor licence (annual)", "$1,000,000", "unchanged"],
        ["Vendor-side algorithm core (Dart port)", "\u2014", "done"],
        ["Customer-side meter simulator (Flutter)", "\u2014", "done"],
        ["NectarAPI-compatible REST server", "\u2014", "done"],
        ["Official STS standard + test vectors (one-time)", "\u2014", "~$2,500"],
        ["Certification fees, HSM, parallel run", "\u2014", "tracked separately"],
        ["Source code we control (vendor + customer)", "no", "yes"],
        ["Cost of evaluating future alternatives", "high", "low"],
    ]
    table(s, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(4.6),
          rows, body_size=14, header_size=15,
          col_widths=[4.5, 2, 3])
    add_text(s, Inches(0.7), Inches(6.3), SW - Inches(1.4), Inches(0.5),
             "Engineering is no longer the cost \u2014 it is now a sunk asset. "
             "What ~$2,500 buys is the official spec to certify it against.",
             size=15, bold=True, color=NAVY)
    return s


# --- Slide 4: What we've built ----------------------------------------
def s_built():
    s = prs.slides.add_slide(BLANK)
    header(s, "What we've built",
           "Three components, one pure-Dart stack, end-to-end working today.")
    rows = [
        ["Component", "Role", "Status"],
        ["nectar_sts_dart  (library)",
         "Pure-Dart implementation of the STS algorithm core: "
         "DKGA-02 / DKGA-04, EA07 (STA), EA09 (DEA / DES), EA11 (MISTY1); "
         "Class 0/1/2 token generators + decoders incl. paired-STA and "
         "4-section MISTY1 KCT; sealed DecodeResult API.",
         "All tests pass"],
        ["nectar_sts_dart  (HTTP server)",
         "NectarAPI-compatible shelf server \u2014 POST /v1/tokens, "
         "POST /v1/tokens/{tokenNo} (decode), /v1/meters registry, "
         "JSON + optional MySQL persistence. Same request shape as the "
         "existing Java tokens-service.",
         "Running"],
        ["nectar_virtual_meter  (Flutter)",
         "Customer-side meter app on Windows / Android / Web / Linux. "
         "Personalize -> punch in 20-digit token -> balance updates. "
         "Handles all 13 ApplyResult variants (credit, replay, KCT staging, "
         "rotation, tamper, Class 2 admin).",
         "Built + tested"],
    ]
    table(s, Inches(0.5), Inches(1.5), SW - Inches(1.0), Inches(4.6),
          rows, body_size=13, header_size=14,
          col_widths=[2.6, 6.6, 1.2])
    add_rect(s, Inches(0.5), Inches(6.3), SW - Inches(1.0), Inches(0.6),
             LIGHT)
    add_text(s, Inches(0.7), Inches(6.42), SW - Inches(1.4), Inches(0.5),
             "Repositories:  dart/nectar_sts_dart  \u00b7  "
             "flutter/nectar_virtual_meter  \u00b7  "
             "java/tokens-service (reference)",
             size=12, color=NAVY, font="Consolas")
    return s


# --- Slide 5: About the STS Association --------------------------------
def s_sts_assn():
    s = prs.slides.add_slide(BLANK)
    header(s, "About the STS Association",
           "Custodian of the global prepayment-electricity standard.")
    bullets(s, Inches(0.9), Inches(1.6), SW - Inches(1.8), Inches(4.4), [
        "STS = Standard Transfer Specification (IEC 62055-41) \u2014 the "
        "global open standard for one-way prepayment electricity tokens.",
        "The STS Association maintains the specification, licenses the "
        "cryptographic algorithms (STA, DEA, MISTY1, DKGA01\u201304), and "
        "certifies vendor implementations.",
        "Every STS-compliant meter \u2014 across vendors, across countries "
        "\u2014 speaks the same 20-digit token format. That is what makes "
        "a Plan B possible.",
        "Only the STS Association can license the standard. A vendor API "
        "into their security module is convenience, not independence "
        "\u2014 pricing and access stay theirs.",
        "Our ~$2,500 buys the official copy of the standard and the test "
        "vectors we will use to certify our pure-Dart implementation.",
    ], size=15)
    add_rect(s, Inches(0.7), Inches(6.0), SW - Inches(1.4), Inches(0.7),
             LIGHT)
    add_text(s, Inches(0.9), Inches(6.13), SW - Inches(1.8), Inches(0.5),
             "Based in South Africa  \u00b7  sts.org.za",
             size=14, color=MUTED)
    return s


# --- Slide 6: IEC 62055-41 alignment ----------------------------------
def s_aligned():
    s = prs.slides.add_slide(BLANK)
    header(s, "IEC 62055-41 ed3.0 alignment \u2014 implemented today",
           "Every row below is shipping in the current Dart stack.")
    rows = [
        ["STS feature", "Status"],
        ["66-bit token framing + 2-bit class transposition", "\u2713"],
        ["Token classes per Table 14 (Class 0 / 1 / 2)", "\u2713"],
        ["Class 0 SubClass 0 \u2014 Transfer Electricity Credit", "\u2713"],
        ["Class 1 SubClasses 0\u20131 \u2014 Initiate Meter Test / Display",
         "\u2713"],
        ["Class 2 register family \u2014 MaxPower / ClearCredit / Tariff / "
         "Tamper / MPPUL", "\u2713"],
        ["Class 2 KCT \u2014 2-section STA + 4-section MISTY1 key change",
         "\u2713"],
        ["DKGA-02 (DES-based) decoder-key derivation", "\u2713"],
        ["DKGA-04 (HMAC-SHA-256, 160-bit VK, MISTY1-capable)", "\u2713"],
        ["EA07 STA  \u00b7  EA09 DEA  \u00b7  EA11 MISTY1 block ciphers",
         "\u2713"],
        ["TID = minutes since BaseDate (1993 / 2014 / 2035)", "\u2713"],
        ["20-digit decimal customer display + CRC", "\u2713"],
        ["Tamper latch + ClearTamperCondition (Class 2/5)", "\u2713"],
        ["Per-meter replay protection (TID window in meter)", "\u2713"],
        ["Vendor-side audit log + meter registry", "\u2713"],
        ["Customer-side meter simulator (Flutter, 4 platforms)", "\u2713"],
    ]
    tbl = table(s, Inches(0.7), Inches(1.5), SW - Inches(1.4), Inches(5.5),
                rows, body_size=12, header_size=13,
                col_widths=[9, 1])
    from pptx.enum.text import PP_ALIGN
    for r_idx in range(1, len(rows)):
        cell = tbl.cell(r_idx, 1)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = TEAL
    return s


# --- Slide 7: What's left before production ---------------------------
def s_swap():
    s = prs.slides.add_slide(BLANK)
    header(s, "What's left before production \u2014 the cert path",
           "The engineering is done. The remaining work is verification + governance.")
    rows = [
        ["Step", "What it is", "Who owns it"],
        ["1.  Buy the official IEC 62055-41 ed3.0 standard + test vectors",
         "~$2,500 one-time. Unlocks bit-exact verification.",
         "Procurement"],
        ["2.  Run KAT (Known-Answer Test) suite against EA07/09/11 + DKGA",
         "Compare our pure-Dart implementations against the official "
         "test vectors. Any mismatch \u2192 fix and re-run.",
         "Engineering"],
        ["3.  STS Association compliance certification",
         "Submit the implementation for STSA review and certification "
         "as a registered Vending System.",
         "Engineering + STSA"],
        ["4.  Vending key management (HSM)",
         "Move the Vending Key off process memory into a hardware HSM "
         "(or PKCS#11 token). Code is already HSM-shaped (see PrismHsm).",
         "Ops + Security"],
        ["5.  Shadow run alongside vendor",
         "1 billing cycle. Every issued token reconciled against the "
         "vendor's output. Customer impact: zero.",
         "Engineering + Ops"],
    ]
    table(s, Inches(0.4), Inches(1.5), SW - Inches(0.8), Inches(4.8),
          rows, body_size=12, header_size=13,
          col_widths=[3.5, 5.5, 1.6])
    add_rect(s, Inches(0.4), Inches(6.5), SW - Inches(0.8), Inches(0.6),
             LIGHT)
    add_text(s, Inches(0.6), Inches(6.62), SW - Inches(1.2), Inches(0.5),
             "Original deck called the cipher a stand-in. As of today it is "
             "the real STA / DEA / MISTY1, pure-Dart \u2014 only "
             "certification remains.",
             size=13, bold=True, color=NAVY)
    return s


# --- Slide 8: Live demo ------------------------------------------------
def s_demo():
    s = prs.slides.add_slide(BLANK)
    header(s, "Live demo \u2014 what you'll see next",
           "Two terminals + the Flutter app  \u00b7  under two minutes")
    bullets(s, Inches(0.9), Inches(1.6), SW - Inches(1.8), Inches(5.4), [
        "Start the Dart vending server  \u2192  POST /v1/tokens accepts "
        "requests",
        "Open the Flutter virtual meter on Windows / Web  \u2192  "
        "personalize from the same vending key + identity",
        "Server issues a 25 kWh credit token  \u2192  20-digit number "
        "printed",
        "Punch the token into the Flutter app  \u2192  balance jumps to "
        "25 kWh, applied-token log records it",
        "Apply the same token again  \u2192  rejected as replay (TID seen)",
        "Trigger tamper sensor in the app  \u2192  Dashboard shows "
        "TAMPER LATCHED",
        "Server issues ClearTamperCondition (Class 2/5)  \u2192  meter "
        "clears it",
        "Server issues paired KCT (STA) or 4-section KCT (MISTY1)  "
        "\u2192  meter stages each section, rotates decoder key on the "
        "final one",
    ], size=16)
    return s


# --- Slide 9: Plan -----------------------------------------------------
def s_plan():
    s = prs.slides.add_slide(BLANK)
    header(s, "Proposed plan \u2014 gated R&D track v2",
           "Vendor contract stays in place until every gate passes.")
    rows = [
        ["#", "Activity", "Duration", "Gate"],
        ["1", "Sponsorship + purchase approval", "1 week",
         "Go / no-go"],
        ["2", "Buy IEC 62055-41 ed3.0 + STS Association test vectors",
         "1\u20132 weeks", "Received"],
        ["3", "Run bit-exact KAT suite against EA07 / EA09 / EA11 + "
              "DKGA-02/04",
         "1\u20132 weeks", "All vectors match"],
        ["4", "Move Vending Key into HSM (PKCS#11 or Prism HSM)",
         "2\u20133 weeks", "HSM signed off"],
        ["5", "STS Association certification",
         "per STSA schedule", "Certified"],
        ["6", "Shadow run alongside vendor (zero customer impact)",
         "1 billing cycle", "100% token match"],
        ["7", "Hand over as standby capability \u2014 ready for renewal "
              "+ contingency",
         "\u2014", "Decision"],
    ]
    table(s, Inches(0.4), Inches(1.6), SW - Inches(0.8), Inches(5.2),
          rows, body_size=13, header_size=14,
          col_widths=[0.3, 5, 1.6, 1.8])
    return s


# --- Slide 10: Risks ---------------------------------------------------
def s_risks():
    s = prs.slides.add_slide(BLANK)
    header(s, "Risk & mitigation")
    rows = [
        ["Risk", "Mitigation"],
        ["Pure-Dart cipher differs from licensed reference",
         "Step 3 \u2014 run the STSA Known-Answer Tests. Any mismatch is "
         "fixed before submission, not in production."],
        ["Certification delays",
         "Vendor stays in place until certification + parallel run pass. "
         "No customer-facing change until both gates are green."],
        ["Vending Key handling",
         "HSM integration is Step 4 of the plan. Code is already "
         "HSM-shaped (VirtualHsm / PrismHsm interfaces ready for swap-in)."],
        ["Single-engineer bus factor",
         "Three documented sub-projects, full test suites, README + "
         "design docs. Onboarding is days, not weeks."],
        ["Customer-side simulator drifts from real meters",
         "Token wire format is fixed by IEC 62055-41 \u2014 every STS "
         "meter consumes the same 20-digit string. The Flutter app is "
         "for our test labs, not customers."],
        ["Vendor lock-in on existing meters",
         "All STS meters speak the same 20-digit format \u2014 no meter "
         "firmware changes required."],
    ]
    table(s, Inches(0.5), Inches(1.5), SW - Inches(1.0), Inches(5.0),
          rows, body_size=13, header_size=14,
          col_widths=[3.6, 6.4])
    return s


# --- Slide 11: Anticipated questions ----------------------------------
def s_qa():
    s = prs.slides.add_slide(BLANK)
    header(s, "Anticipated questions")
    qa = [
        ("Is $2,500 the total cost?",
         "No. It buys the official IEC 62055-41 ed3.0 standard and the "
         "STSA test vectors. Certification fees, HSM and shadow-run "
         "effort are tracked separately."),
        ("Is the cipher still a stand-in?",
         "No \u2014 since the 5-Jun deck we now have pure-Dart "
         "implementations of STA (EA07), DEA (EA09) and MISTY1 (EA11), "
         "plus DKGA-02 and DKGA-04. What's missing is third-party "
         "test-vector verification + STSA cert."),
        ("How is this different from the original electricity-token-generator?",
         "That project was a teaching analogue (Feistel + HMAC stand-ins). "
         "nectar_sts_dart is a clean-room implementation of the real STS "
         "algorithms, ported from the Java tokens-service we already run."),
        ("Why a Flutter app?",
         "Customer-side meter simulator for our test labs and field "
         "engineering. Same code on Windows / Android / Web / Linux. "
         "Provisioning + token-apply round-trip is in the docs."),
        ("What if Kelly leaves?",
         "All three components have tests, READMEs, and a docs/ folder. "
         "Any Dart / Flutter engineer can take over in days."),
    ]
    y = Inches(1.3)
    for q, a in qa:
        add_text(s, Inches(0.7), y, SW - Inches(1.4), Inches(0.4),
                 f"Q.  {q}", size=14, bold=True, color=NAVY)
        add_text(s, Inches(0.9), y + Inches(0.36),
                 SW - Inches(1.6), Inches(0.6),
                 f"A.  {a}", size=12, color=TEXT)
        y += Inches(1.05)
    return s


# --- Slide 12: The decision -------------------------------------------
def s_decision():
    s = prs.slides.add_slide(BLANK)
    header(s, "The decision being asked for",
           "No vendor contract is touched.")
    items = [
        ("1", "R&D sponsorship (renewed)",
         "Continue the in-house STS track with a named technical lead. "
         "Engineering scope already delivered."),
        ("2", "Approval to buy the official STS standard + test vectors",
         "~$2,500, one-time. Unlocks bit-exact KAT verification."),
        ("3", "Review checkpoint after the shadow run",
         "Decision: hold as standby, expand to active vending, or shelve."),
    ]
    y = Inches(1.6)
    for num, head, body in items:
        add_rect(s, Inches(0.7), y, Inches(0.8), Inches(1.2), ACCENT)
        add_text(s, Inches(0.7), y + Inches(0.25), Inches(0.8),
                 Inches(0.8), num, size=36, bold=True, color=WHITE,
                 align="center")
        add_text(s, Inches(1.7), y + Inches(0.05),
                 SW - Inches(2.4), Inches(0.5),
                 head, size=20, bold=True, color=NAVY)
        add_text(s, Inches(1.7), y + Inches(0.55),
                 SW - Inches(2.4), Inches(0.7),
                 body, size=14, color=TEXT)
        y += Inches(1.55)
    return s


# --- Slide 13: Glossary -----------------------------------------------
def s_glossary():
    s = prs.slides.add_slide(BLANK)
    header(s, "Glossary  \u00b7  acronyms used in this deck",
           "For reference during questions.")
    rows = [
        ["Term", "Expansion"],
        ["API", "Application Programming Interface"],
        ["DEA  (EA09)",
         "Data Encryption Algorithm \u2014 DES-based STS block cipher"],
        ["DES",
         "Data Encryption Standard (the underlying NIST cipher used by DEA)"],
        ["DKGA01\u201304",
         "Decoder Key Generation Algorithms \u2014 key-derivation functions; "
         "DKGA04 added in ed3.0 for 160-bit Vending Keys"],
        ["EA07 / EA09 / EA11",
         "STS Encryption Algorithm identifiers  (= STA / DEA / MISTY1)"],
        ["HMAC-SHA-256",
         "Hash-based MAC using SHA-256 \u2014 used by DKGA-04"],
        ["HSM", "Hardware Security Module (tamper-resistant key store)"],
        ["IAIN",
         "Individual Account Identification Number \u2014 meter's account "
         "portion of the PAN"],
        ["IEC", "International Electrotechnical Commission"],
        ["IIN",
         "Issuer Identification Number \u2014 utility identifier portion "
         "of the PAN"],
        ["KAT",
         "Known-Answer Test \u2014 vendor-supplied test vectors for "
         "cipher verification"],
        ["KCT",
         "Key Change Token \u2014 2 (STA) or 4 (MISTY1) coordinated "
         "tokens that rotate the decoder key"],
        ["KMS", "Key Management System"],
        ["KRN",
         "Key Revision Number \u2014 which generation of derived key the "
         "token used"],
        ["kWh", "Kilowatt-hour (the credit unit a Class 0/0 token carries)"],
        ["MISTY1  (EA11)",
         "Mitsubishi Improved Security cipher licensed for STS use"],
        ["MPPUL",
         "Maximum Phase Power Unbalance Limit \u2014 Class 2/6 management value"],
        ["PAN",
         "Primary Account Number \u2014 IIN + IAIN, the per-meter identity"],
        ["SGC",
         "Supply Group Code \u2014 operator / region identifier in the "
         "key hierarchy"],
        ["STA  (EA07)",
         "STS Token Algorithm \u2014 the proprietary STS block cipher"],
        ["STS", "Standard Transfer Specification  (IEC 62055-41)"],
        ["STSA",
         "STS Association \u2014 custodian and licensor of the STS standard"],
        ["TI",
         "Tariff Index \u2014 2-digit selector into the meter's tariff table"],
        ["TID",
         "Token Identifier \u2014 minutes since the BaseDate, ensures "
         "uniqueness"],
        ["UTC", "Coordinated Universal Time"],
        ["VK", "Vending Key \u2014 root key at the utility / vending level"],
    ]
    table(s, Inches(0.4), Inches(1.4), SW - Inches(0.8), Inches(5.7),
          rows, body_size=10, header_size=12,
          col_widths=[2, 8])
    return s


# --- Slide 14: Closing ------------------------------------------------
def s_close():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_text(s, Inches(0.8), Inches(2.6), SW - Inches(1.6), Inches(1.2),
             "Thank you.", size=60, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.9), SW - Inches(1.6), Inches(0.6),
             "Questions \u2014 and a live demo on request.",
             size=22, color=LIGHT)
    add_rect(s, Inches(0.8), Inches(5.0), Inches(4), Inches(0.04), ACCENT)
    add_text(s, Inches(0.8), Inches(5.2), SW - Inches(1.6), Inches(0.4),
             "Kelly Kinyama  \u00b7  11 June 2026", size=14, color=LIGHT)
    return s


builders = [
    s_title, s_ask, s_numbers, s_built, s_sts_assn, s_aligned, s_swap,
    s_demo, s_plan, s_risks, s_qa, s_decision, s_glossary, s_close,
]
total = len(builders)
for i, b in enumerate(builders, 1):
    slide = b()
    if i not in (1, total):  # no footer on title / closing
        footer(slide, i, total)

prs.save(OUT)
print(f"Wrote {OUT}  ({OUT.stat().st_size:,} bytes, {total} slides)")
prs.save(OUT_ALT)
print(f"Wrote {OUT_ALT}  ({OUT_ALT.stat().st_size:,} bytes, {total} slides)")
